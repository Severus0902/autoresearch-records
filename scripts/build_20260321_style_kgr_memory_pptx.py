# -*- coding: utf-8 -*-
"""Extend the 20260321 KGR progress deck with experiments and memory benchmarks.

The original deck is used as the visual template and is never overwritten.
Run from the repository root with Python 3.8:

    $env:PYTHONPATH='.tmp/pptx_native_py38'
    D:\Python38\python.exe scripts/build_20260321_style_kgr_memory_pptx.py
"""

from __future__ import annotations

import sys
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
NATIVE_DEPS = ROOT / ".tmp" / "pptx_native_py38"
PURE_DEPS = ROOT / ".tmp" / "pptx_deps"
if NATIVE_DEPS.exists():
    sys.path.insert(0, str(NATIVE_DEPS))
if PURE_DEPS.exists():
    # Keep the system typing_extensions ahead of the newer wheel copied here.
    sys.path.append(str(PURE_DEPS))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SRC = ROOT / "docs" / "slides" / "20260321.pptx"
OUT = ROOT / "docs" / "slides" / "2026-09-03-kgr-agent-memory-progress.pptx"
FRAMEWORK = ROOT / "docs" / "figures" / "2026-09-03-memreadybench-stage1-benchmark-framework.png"

SW = 13.333
SH = 7.5
FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"

NAVY = RGBColor(35, 79, 126)
BLUE = RGBColor(45, 106, 170)
BLUE_LIGHT = RGBColor(231, 241, 251)
GOLD = RGBColor(255, 214, 92)
ORANGE = RGBColor(244, 160, 0)
INK = RGBColor(27, 31, 36)
MUTED = RGBColor(92, 101, 112)
LINE = RGBColor(205, 211, 219)
LIGHT = RGBColor(247, 248, 250)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(39, 119, 73)
GREEN_LIGHT = RGBColor(231, 246, 237)
RED = RGBColor(183, 56, 56)
RED_LIGHT = RGBColor(252, 235, 235)
PURPLE = RGBColor(104, 75, 150)
PURPLE_LIGHT = RGBColor(241, 236, 249)
TEAL = RGBColor(31, 126, 128)
TEAL_LIGHT = RGBColor(230, 246, 246)


def set_text_style(run, size, color=INK, bold=False, font=FONT_CN):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
    font=FONT_CN,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    lines = str(text).split("\n")
    for idx, line_text in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line_text
        set_text_style(run, size, color=color, bold=bold, font=font)
    return box


def add_rect(slide, x, y, w, h, *, fill=WHITE, line=LINE, width=1.0, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(width)
    return shp


def add_line(slide, x1, y1, x2, y2, *, color=LINE, width=1.5, arrow=False):
    shp = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    shp.line.color.rgb = color
    shp.line.width = Pt(width)
    if arrow:
        shp.line.end_arrowhead = True
    return shp


def add_title(slide, title):
    return add_text(slide, title, 1.39, 0.18, 10.5, 0.58, size=27, bold=True, valign=MSO_ANCHOR.MIDDLE)


def add_slide_number(slide, number):
    add_text(slide, str(number), 12.20, 6.95, 0.48, 0.24, size=8, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT_EN)


def add_citation(slide, citation, *, width=9.8):
    add_rect(slide, 0.42, 6.78, width, 0.23, fill=WHITE, line=WHITE, width=0)
    add_text(slide, citation, 0.46, 6.79, width - 0.08, 0.18, size=7.2, color=MUTED, font=FONT_EN)


def add_body_slide(prs, title, number, citation=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title)
    add_slide_number(slide, number)
    if citation:
        add_citation(slide, citation)
    return slide


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    for placeholder in slide.placeholders:
        placeholder.text = ""
    add_text(
        slide,
        title,
        6.24,
        3.14,
        6.45,
        0.88,
        size=40 if len(title) > 12 else 46,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return slide


def add_bullet_list(slide, items, x, y, w, h, *, size=18, color=INK, gap=8, bullet_color=NAVY):
    total = len(items)
    row_h = h / max(total, 1)
    for idx, item in enumerate(items):
        cy = y + idx * row_h
        add_text(slide, "▸", x, cy + 0.01, 0.28, row_h - 0.02, size=size, color=bullet_color, bold=True)
        add_text(slide, item, x + 0.30, cy, w - 0.30, row_h - 0.02, size=size, color=color, line_spacing=1.06)


def add_metric_card(slide, x, y, w, h, value, label, note, *, accent=NAVY, fill=LIGHT):
    add_rect(slide, x, y, w, h, fill=fill, line=accent, width=1.2)
    add_text(slide, value, x + 0.15, y + 0.12, w - 0.30, 0.58, size=30, color=accent, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(slide, label, x + 0.15, y + 0.76, w - 0.30, 0.28, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, note, x + 0.16, y + 1.08, w - 0.32, h - 1.18, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)


def add_step(slide, x, y, w, h, number, title, body, *, accent=NAVY, fill=LIGHT):
    add_rect(slide, x, y, w, h, fill=fill, line=accent, width=1.2)
    add_rect(slide, x, y, 0.48, h, fill=accent, line=accent, width=0)
    add_text(slide, str(number), x + 0.04, y + 0.25, 0.39, 0.45, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(slide, title, x + 0.62, y + 0.14, w - 0.76, 0.32, size=15, color=accent, bold=True)
    add_text(slide, body, x + 0.62, y + 0.55, w - 0.76, h - 0.62, size=11, color=INK, line_spacing=1.03)


def add_tag(slide, text, x, y, w, *, accent=NAVY, fill=BLUE_LIGHT):
    add_rect(slide, x, y, w, 0.36, fill=fill, line=accent, width=0.8, radius=True)
    add_text(slide, text, x + 0.04, y + 0.07, w - 0.08, 0.20, size=10, color=accent, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_table(slide, rows, columns, x, y, w, h, *, widths=None, header_fill=NAVY, font_size=11):
    table = slide.shapes.add_table(len(rows), len(columns), Inches(x), Inches(y), Inches(w), Inches(h)).table
    if widths:
        for col, width in zip(table.columns, widths):
            col.width = Inches(width)
    data = [columns] + rows[1:] if rows and rows[0] == columns else [columns] + rows
    if len(data) != len(rows):
        # Caller normally passes only body rows; resize is fixed, so rebuild defensively.
        raise ValueError("Rows must include enough entries for the table dimensions")
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else (WHITE if r % 2 else LIGHT)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_text_style(run, font_size, color=WHITE if r == 0 else INK, bold=(r == 0))
    return table


def add_bar(slide, label, value, max_value, x, y, w, *, color=NAVY, value_fmt="{:.1%}"):
    add_text(slide, label, x, y, 1.85, 0.28, size=11.5, color=INK)
    add_rect(slide, x + 1.88, y + 0.04, w - 2.72, 0.20, fill=RGBColor(230, 233, 238), line=RGBColor(230, 233, 238), width=0)
    bar_w = max(0.02, (w - 2.72) * value / max_value)
    add_rect(slide, x + 1.88, y + 0.04, bar_w, 0.20, fill=color, line=color, width=0)
    add_text(slide, value_fmt.format(value), x + w - 0.78, y, 0.75, 0.28, size=11, color=color, bold=True, align=PP_ALIGN.RIGHT, font=FONT_EN)


def style_existing_text(shape, size, *, color=INK, bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in tf.paragraphs:
        p.alignment = align
        for run in p.runs:
            set_text_style(run, size, color=color, bold=bold)


def add_reference_columns(slide, left_refs, right_refs):
    for x, refs, start in ((0.58, left_refs, 1), (6.70, right_refs, len(left_refs) + 1)):
        y = 1.43
        for idx, ref in enumerate(refs, start):
            add_text(slide, f"[{idx}]", x, y, 0.34, 0.25, size=10, color=NAVY, bold=True, font=FONT_EN)
            add_text(slide, ref, x + 0.38, y - 0.01, 5.72, 0.72, size=9.3, color=INK, font=FONT_EN, line_spacing=1.0)
            y += 0.98


def main():
    if not SRC.exists():
        raise FileNotFoundError(SRC)
    if not FRAMEWORK.exists():
        raise FileNotFoundError(FRAMEWORK)

    prs = Presentation(str(SRC))
    closing_sld_id = prs.slides._sldIdLst[-1]

    # Refresh the cover while preserving the original template geometry.
    cover = prs.slides[0]
    cover.shapes[0].text = "20260903进展汇报"
    style_existing_text(cover.shapes[0], 44, color=INK, align=PP_ALIGN.CENTER)
    cover.shapes[1].text = "LLM-based KGR验证与Agent Memory Benchmark\n汇报人：魏熙润"
    style_existing_text(cover.shapes[1], 24, color=INK, align=PP_ALIGN.CENTER)

    # Page-level source notes for the papers already present in the March deck.
    add_citation(prs.slides[4], "Source: Liu et al., PathMind, AAAI 2026, Fig. 1.", width=5.1)
    add_citation(prs.slides[5], "Source: Liu et al., PathMind, AAAI 2026; Pan et al., IEEE TKDE 2024.", width=6.2)
    add_citation(prs.slides[6], "Source: Sun et al., ToG, ICLR 2024; Liu et al., PathMind, AAAI 2026.", width=6.2)
    add_citation(prs.slides[7], "Source: Liu et al., PathMind, AAAI 2026.", width=5.1)
    add_citation(prs.slides[9], "Source: Jiang et al., KG-Agent, ACL 2025.", width=5.0)
    add_citation(prs.slides[11], "Source: Wu et al., Agentic Reasoning, ACL 2025.", width=5.5)
    add_citation(prs.slides[12], "Source: Wu et al., Agentic Reasoning, ACL 2025.", width=5.5)

    page = 15

    add_section_slide(prs, "KGR实验验证")

    slide = add_body_slide(
        prs,
        "EoG / BoG：创新集中在探索策略，而非子图检索",
        page,
        "Yan et al., Explore-on-Graph, ICLR 2026; Zhang et al., Backjump-on-Graph, ICML 2026 (PMLR 306).",
    )
    add_rect(slide, 0.58, 1.18, 5.92, 4.65, fill=BLUE_LIGHT, line=BLUE, width=1.2)
    add_text(slide, "Explore-on-Graph (EoG)", 0.92, 1.47, 5.24, 0.38, size=20, color=BLUE, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_bullet_list(
        slide,
        [
            "动机：规则约束或固定示范路径限制 OOD 泛化与新路径探索。",
            "Stage 1：用长 CoT / 图探索轨迹进行 SFT 冷启动。",
            "Stage 2：GRPO + final-answer correctness，并用 path information 细化过程奖励。",
            "重点：鼓励多样、有效的自主探索，减少无意义路径。",
        ],
        0.92,
        2.08,
        5.20,
        2.92,
        size=14.3,
        bullet_color=BLUE,
    )
    add_tag(slide, "explore broader", 2.03, 5.20, 2.96, accent=BLUE, fill=WHITE)
    add_rect(slide, 6.82, 1.18, 5.92, 4.65, fill=PURPLE_LIGHT, line=PURPLE, width=1.2)
    add_text(slide, "Backjump-on-Graph (BoG)", 7.16, 1.47, 5.24, 0.38, size=20, color=PURPLE, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_bullet_list(
        slide,
        [
            "动机：query想象的关系路径与真实 KG 拓扑不匹配，Agent 容易进入 dead end。",
            "Stage 1：合成含 Forward / Backjump / Yield / Halt 的轨迹进行 SFT。",
            "Stage 2：RL优化回跳时机与落点；混合奖励抑制冗余转移并奖励正确答案。",
            "重点：从历史状态非局部回跳，再探索替代分支。",
        ],
        7.16,
        2.08,
        5.20,
        2.92,
        size=14.3,
        bullet_color=PURPLE,
    )
    add_tag(slide, "recover from dead ends", 8.10, 5.20, 3.36, accent=PURPLE, fill=WHITE)
    add_text(slide, "共同空缺：仍未显式把同一步的候选 action/path 竞争建模为 ranking problem，也未解决 memory 何时有益或有害。", 0.82, 6.12, 11.72, 0.43, size=17, color=RED, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    slide = add_body_slide(
        prs,
        "KGR最小闭环：从候选动作到记忆门控",
        page,
        "Internal experiment logs: docs/results/2026-09-01 to 2026-09-02.",
    )
    steps = [
        ("Query + topic entity", "WebQSP问题与中心实体"),
        ("Local subgraph", "Freebase一跳关系候选"),
        ("Action ranking", "rule / memory / LR ranker"),
        ("Qwen3-0.6B", "SFT受控输出下一跳"),
        ("Memory gate", "决定是否注入经验"),
    ]
    x_positions = [0.55, 3.02, 5.49, 7.96, 10.43]
    for idx, ((title, body), x) in enumerate(zip(steps, x_positions), 1):
        add_step(slide, x, 1.35, 2.28, 1.55, idx, title, body, accent=NAVY if idx < 4 else TEAL, fill=BLUE_LIGHT if idx < 4 else TEAL_LIGHT)
        if idx < len(steps):
            add_line(slide, x + 2.28, 2.12, x + 2.46, 2.12, color=ORANGE, width=2.3, arrow=True)
    add_rect(slide, 0.76, 3.46, 11.80, 2.20, fill=LIGHT, line=LINE, width=1.0)
    add_text(slide, "验证目标", 1.03, 3.72, 1.35, 0.34, size=18, color=NAVY, bold=True)
    add_bullet_list(
        slide,
        [
            "先验证候选空间中是否存在可学习的排序信号，再验证小模型能否稳定输出动作。",
            "比较 no-memory / random-memory / verified-memory，判断跨样本经验是否真正改变策略。",
            "训练轻量 gate，检验“何时用 memory”能否接近 oracle，而非默认把 memory 拼进 prompt。",
        ],
        2.42,
        3.67,
        9.58,
        1.58,
        size=16,
    )
    add_text(slide, "结论：闭环跑通，但性能瓶颈依次暴露在检索覆盖、记忆门控和排序到策略的转化。", 1.03, 5.92, 11.25, 0.40, size=18, color=RED, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    slide = add_body_slide(
        prs,
        "瓶颈首先出现在子图构建",
        page,
        "Stage 1 logs: WebQSP eval100/train500; answer visibility is computed in the sampled local subgraph.",
    )
    add_metric_card(slide, 0.62, 1.30, 2.65, 2.00, "0.15", "初始候选关系召回", "SELECT ?r ?o LIMIT 50 被元数据关系和重复对象占满", accent=RED, fill=RED_LIGHT)
    add_text(slide, "→", 3.42, 1.83, 0.52, 0.48, size=32, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_metric_card(slide, 4.05, 1.30, 2.65, 2.00, "0.86", "修复后候选召回", "关系级采样 + namespace过滤；eval100", accent=GREEN, fill=GREEN_LIGHT)
    add_metric_card(slide, 7.22, 1.30, 2.65, 2.00, "0.94", "边召回", "目标关系相关边进入候选子图；eval100", accent=NAVY, fill=BLUE_LIGHT)
    add_metric_card(slide, 10.05, 1.30, 2.65, 2.00, "0.62", "答案可见率", "即使关系召回较高，答案实体仍只在62%子图中可见", accent=PURPLE, fill=PURPLE_LIGHT)
    add_rect(slide, 0.75, 3.78, 5.65, 2.05, fill=RED_LIGHT, line=RED, width=1.1)
    add_text(slide, "错误认识", 1.00, 4.02, 1.32, 0.30, size=17, color=RED, bold=True)
    add_text(slide, "只要后续推理策略更强，就能解决 KGQA。", 1.00, 4.50, 4.95, 0.40, size=18, bold=True)
    add_text(slide, "实际：检索漏召回会直接截断答案路径，policy 无法恢复不存在的证据。", 1.00, 5.12, 4.95, 0.45, size=14, color=MUTED)
    add_rect(slide, 6.85, 3.78, 5.72, 2.05, fill=GREEN_LIGHT, line=GREEN, width=1.1)
    add_text(slide, "得到的设计约束", 7.10, 4.02, 2.10, 0.30, size=17, color=GREEN, bold=True)
    add_bullet_list(slide, ["必须报告 candidate recall / answer visibility。", "推理改进必须与检索上限分开归因。", "gold subgraph 只能作为 oracle，不是公平主设定。"], 7.08, 4.43, 5.00, 1.20, size=13.5, bullet_color=GREEN)
    page += 1

    slide = add_body_slide(
        prs,
        "0.6B模型：SFT能学会动作格式，memory增益很小",
        page,
        "Stage 6: Qwen3-0.6B, 256 training samples, 30 SFT steps; ablation evaluated on 60 examples.",
    )
    add_metric_card(slide, 0.65, 1.30, 2.55, 1.78, "0%", "SFT前准确率", "invalid output rate = 100%", accent=RED, fill=RED_LIGHT)
    add_text(slide, "→", 3.36, 1.79, 0.48, 0.42, size=30, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_metric_card(slide, 4.00, 1.30, 2.55, 1.78, "70%", "SFT后准确率", "invalid output rate = 0%", accent=GREEN, fill=GREEN_LIGHT)
    add_text(slide, "说明小模型可用于验证受控 action policy，但不能把格式学习误判为 memory 有效。", 0.76, 3.32, 5.72, 0.80, size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 6.95, 1.22, 5.65, 4.62, fill=LIGHT, line=LINE, width=1.0)
    add_text(slide, "Memory ablation / Acc@1", 7.28, 1.48, 4.98, 0.32, size=17, color=NAVY, bold=True)
    add_bar(slide, "No memory", 0.7167, 0.80, 7.24, 2.14, 4.95, color=MUTED)
    add_bar(slide, "Random memory", 0.7167, 0.80, 7.24, 2.83, 4.95, color=PURPLE)
    add_bar(slide, "Verified memory", 0.7333, 0.80, 7.24, 3.52, 4.95, color=GREEN)
    add_text(slide, "+1 / 60", 7.30, 4.38, 4.85, 0.42, size=26, color=GREEN, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(slide, "verified memory 相比 no-memory 只多答对 1 个样本", 7.30, 4.92, 4.85, 0.50, size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "结论：直接拼接 memory 不是稳定创新点，真正的问题是选择、抑制与更新。", 0.82, 5.98, 11.65, 0.42, size=18, color=RED, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    slide = add_body_slide(
        prs,
        "Memory不是单调增益：是否使用比有没有更关键",
        page,
        "Stage 7 memory-gate validation on the same 60-example split.",
    )
    add_rect(slide, 0.62, 1.26, 4.00, 4.67, fill=LIGHT, line=LINE, width=1.0)
    add_text(slide, "逐样本因果现象", 0.92, 1.55, 3.38, 0.34, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_metric_card(slide, 0.96, 2.13, 1.42, 1.45, "6", "helped", "memory纠正", accent=GREEN, fill=GREEN_LIGHT)
    add_metric_card(slide, 2.84, 2.13, 1.42, 1.45, "5", "hurt", "memory误导", accent=RED, fill=RED_LIGHT)
    add_metric_card(slide, 0.96, 3.89, 1.42, 1.45, "38", "both right", "是否注入均正确", accent=NAVY, fill=BLUE_LIGHT)
    add_metric_card(slide, 2.84, 3.89, 1.42, 1.45, "11", "both wrong", "候选/能力上限", accent=MUTED, fill=WHITE)
    add_rect(slide, 4.98, 1.26, 7.68, 4.67, fill=WHITE, line=LINE, width=1.0)
    add_text(slide, "Gate策略 / Acc@1", 5.30, 1.55, 6.98, 0.34, size=19, color=NAVY, bold=True)
    add_bar(slide, "Always verified", 0.7333, 0.90, 5.28, 2.15, 6.90, color=MUTED)
    add_bar(slide, "Heuristic gate", 0.7500, 0.90, 5.28, 2.82, 6.90, color=BLUE)
    add_bar(slide, "GRM-lite (LOOCV)", 0.6833, 0.90, 5.28, 3.49, 6.90, color=RED)
    add_bar(slide, "Oracle: no/verified", 0.8167, 0.90, 5.28, 4.16, 6.90, color=GREEN)
    add_bar(slide, "Oracle: all three", 0.8667, 0.90, 5.28, 4.83, 6.90, color=TEAL)
    add_text(slide, "Oracle gap 说明“选择性使用”有空间；GRM-lite 失败说明样本与状态表征不足。", 1.03, 6.02, 11.25, 0.42, size=17.5, color=RED, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    slide = add_body_slide(
        prs,
        "排序奖励：pairwise学会比较，但尚未转化为Top-1",
        page,
        "Stage 8 minimal loop: WebQSP first-hop candidate ranking, 100-example evaluation.",
    )
    columns = ["Method", "Top-1", "MRR", "R@3", "R@5", "PrefAcc"]
    body = [
        ["rule", ".10", ".193", ".20", ".30", ".494"],
        ["memory", ".49", ".606", ".71", ".73", ".821"],
        ["pointwise LR", ".47", ".598", ".72", ".75", ".881"],
        ["pairwise LR", ".40", ".553", ".70", ".75", ".911"],
        ["listwise RRF", ".33", ".485", ".59", ".73", ".842"],
    ]
    # add_table expects header plus body in a fixed-size matrix.
    table_shape = slide.shapes.add_table(6, 6, Inches(0.62), Inches(1.32), Inches(7.38), Inches(4.55))
    table = table_shape.table
    widths = [2.05, 1.04, 1.04, 1.04, 1.04, 1.17]
    for col, width in zip(table.columns, widths):
        col.width = Inches(width)
    for r, row in enumerate([columns] + body):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = NAVY
            elif r == 4:
                cell.fill.fore_color.rgb = GREEN_LIGHT
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_text_style(run, 12, color=WHITE if r == 0 else INK, bold=(r == 0 or r == 4), font=FONT_EN if c else FONT_CN)
    add_rect(slide, 8.35, 1.32, 4.30, 2.02, fill=GREEN_LIGHT, line=GREEN, width=1.2)
    add_text(slide, ".911", 8.68, 1.58, 3.62, 0.64, size=34, color=GREEN, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(slide, "Pairwise preference accuracy", 8.68, 2.35, 3.62, 0.32, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(slide, "能区分 hard negative", 8.68, 2.76, 3.62, 0.28, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.35, 3.70, 4.30, 2.17, fill=RED_LIGHT, line=RED, width=1.2)
    add_text(slide, "≠  Top-1提升", 8.68, 4.02, 3.62, 0.48, size=24, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Pairwise Top-1 = .40，低于 memory heuristic = .49；简单 listwise 融合继续下降。", 8.72, 4.71, 3.54, 0.78, size=13.5, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, "排序信号适合作为过程奖励/诊断信号，但还需要 policy coupling 才能改进最终答案。", 0.86, 6.08, 11.55, 0.38, size=17.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    slide = add_body_slide(
        prs,
        "KGR实验给出的真实结论",
        page,
        "Scope caveat: current experiments are WebQSP first-hop pilots, not official multi-hop CWQ evaluation.",
    )
    conclusions = [
        ("1", "Retrieval ceiling", "答案不可见时，后续RL/GRPO无法补回证据；必须拆分检索与推理贡献。", RED, RED_LIGHT),
        ("2", "Non-monotonic memory", "memory 对不同样本既有帮助也会误导；平均分掩盖了 use/suppress 决策。", PURPLE, PURPLE_LIGHT),
        ("3", "Gate learning is the bottleneck", "oracle 0.8167，但轻量GRM只有0.6833；缺少可诊断的memory state supervision。", GREEN, GREEN_LIGHT),
        ("4", "Ranking is not policy", "pairwise比较正确并不等于Top-1与答案正确；需要闭环状态、动作与后果。", NAVY, BLUE_LIGHT),
    ]
    coords = [(0.62, 1.28), (6.76, 1.28), (0.62, 3.75), (6.76, 3.75)]
    for (num, title, body, accent, fill), (x, y) in zip(conclusions, coords):
        add_rect(slide, x, y, 5.88, 2.08, fill=fill, line=accent, width=1.2)
        add_text(slide, num, x + 0.22, y + 0.18, 0.48, 0.48, size=25, color=accent, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(slide, title, x + 0.84, y + 0.19, 4.60, 0.33, size=18, color=accent, bold=True, font=FONT_EN)
        add_text(slide, body, x + 0.84, y + 0.73, 4.58, 0.92, size=14, color=INK, line_spacing=1.05)
    add_text(slide, "研究转向：不再把 memory 当额外文本，而是把“状态、使用、抑制、修复”本身变成可测对象。", 0.84, 6.14, 11.60, 0.38, size=18, color=RED, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    add_section_slide(prs, "Agent Memory Benchmarks")

    slide = add_body_slide(
        prs,
        "Agent Memory：评测对象不是长上下文本身",
        page,
        "Conceptual framing: Hu et al., Memory in the Age of AI Agents, arXiv:2512.13564.",
    )
    add_text(slide, "Agent memory = 跨交互持续存在、可被写入/检索/使用/更新，并影响未来行为的认知状态。", 0.72, 1.15, 11.90, 0.55, size=21, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    phases = [
        ("FORM", "形成", "从history / tool / user feedback\n抽取事实、经验与规则", BLUE, BLUE_LIGHT),
        ("STORE", "持久化", "跨session保存内容、时间、\n来源与authority", PURPLE, PURPLE_LIGHT),
        ("RETRIEVE", "检索", "按当前目标召回候选记忆，\n处理更新与冲突", TEAL, TEAL_LIGHT),
        ("USE", "调控", "选择、抑制、澄清或验证，\n再决定是否行动", GREEN, GREEN_LIGHT),
        ("EVOLVE", "演化", "根据世界反馈修正旧记忆，\n支持未来session", ORANGE, RGBColor(255, 245, 220)),
    ]
    x_positions = [0.48, 3.02, 5.56, 8.10, 10.64]
    for idx, ((tag, title, body, accent, fill), x) in enumerate(zip(phases, x_positions)):
        add_rect(slide, x, 2.08, 2.20, 2.55, fill=fill, line=accent, width=1.2)
        add_tag(slide, tag, x + 0.38, 2.27, 1.44, accent=accent, fill=WHITE)
        add_text(slide, title, x + 0.18, 2.82, 1.84, 0.36, size=20, color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.18, 3.42, 1.84, 0.88, size=12.2, color=INK, align=PP_ALIGN.CENTER)
        if idx < 4:
            add_line(slide, x + 2.20, 3.33, x + 2.50, 3.33, color=NAVY, width=1.8, arrow=True)
    add_rect(slide, 0.95, 5.05, 11.44, 0.93, fill=LIGHT, line=LINE, width=1.0)
    add_text(slide, "区别", 1.17, 5.30, 0.75, 0.30, size=16, color=NAVY, bold=True)
    add_text(slide, "RAG主要回答“从外部知识库检索什么”；长上下文主要回答“窗口内能放多少”；Agent Memory还必须回答“哪些状态可被未来行为信任、何时应被抑制或修复”。", 2.02, 5.22, 9.98, 0.48, size=15.3, color=INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    page += 1

    slide = add_body_slide(
        prs,
        "经典Memory Benchmark的演进",
        page,
        "Yang et al., NeurIPS 2024; Maharana et al., ACL 2024; Wu et al., ICLR 2025; Wu et al., ACL 2026; Shen et al., ACL 2026.",
    )
    add_line(slide, 1.05, 3.28, 12.20, 3.28, color=NAVY, width=2.0, arrow=True)
    timeline = [
        ("2024", "CRAG", "动态/长尾事实QA\nRAG可靠性", BLUE, 1.10),
        ("2024", "LoCoMo", "跨session QA\n事件总结", PURPLE, 3.35),
        ("2025", "LongMemEval", "检索 + 推理\n更新 + abstention", TEAL, 5.65),
        ("2026", "StratMem-Bench", "must / nice / irr\n战略使用", GREEN, 8.05),
        ("2026", "Mem2ActBench", "tool selection\nparameter grounding", ORANGE, 10.55),
    ]
    for year, name, body, accent, x in timeline:
        add_rect(slide, x, 1.38, 1.84, 1.32, fill=WHITE, line=accent, width=1.2)
        add_text(slide, year, x + 0.12, 1.52, 1.60, 0.24, size=11, color=accent, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(slide, name, x + 0.10, 1.86, 1.64, 0.37, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(slide, body, x + 0.10, 2.24, 1.64, 0.42, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
        add_line(slide, x + 0.92, 2.70, x + 0.92, 3.28, color=accent, width=1.2)
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.79), Inches(3.14), Inches(0.26), Inches(0.26))
        circle.fill.solid(); circle.fill.fore_color.rgb = accent; circle.line.color.rgb = accent
    add_rect(slide, 1.18, 4.17, 10.95, 1.22, fill=LIGHT, line=LINE, width=1.0)
    add_text(slide, "演进主线", 1.48, 4.50, 1.25, 0.30, size=17, color=NAVY, bold=True)
    add_text(slide, "检索到事实  →  跨会话记忆  →  多类型能力  →  战略选择  →  由记忆驱动工具行动", 2.92, 4.43, 8.78, 0.44, size=19, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "仍缺少统一的、受控的 memory-state intervention：同一个任务仅改变 stale/conflicting/incomplete/verified memory 后，Agent 是否会改变行为？", 1.16, 5.75, 11.00, 0.56, size=16, color=RED, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    slide = add_body_slide(
        prs,
        "经典工作对照：测到了什么，还没测什么",
        page,
        "The first four works are explicitly discussed in the StratMem-Bench introduction/related work.",
    )
    headers = ["Benchmark", "Memory source", "Core task", "Primary signal", "Remaining gap"]
    rows = [
        ["CRAG", "Web / KG API", "factual QA", "accuracy / hallucination", "非持久memory lifecycle"],
        ["LoCoMo", "长对话事件", "QA / summary", "long-range recall", "不控制memory状态"],
        ["LongMemEval", "历史会话", "5类memory能力", "QA accuracy", "以回答为终点"],
        ["StratMem", "must/nice/irr池", "条件响应生成", "使用/抑制/整合", "不执行环境动作"],
        ["Mem2Act", "长期任务链", "tool + parameters", "tool-call accuracy", "缺少verify/repair闭环"],
    ]
    shape = slide.shapes.add_table(6, 5, Inches(0.44), Inches(1.18), Inches(12.45), Inches(4.98))
    table = shape.table
    widths = [1.42, 2.05, 2.20, 2.45, 4.33]
    for col, width in zip(table.columns, widths):
        col.width = Inches(width)
    for r, row in enumerate([headers] + rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else (WHITE if r % 2 else LIGHT)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.05); cell.margin_right = Inches(0.05)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_text_style(run, 11.2, color=WHITE if r == 0 else INK, bold=(r == 0), font=FONT_EN if c == 0 else FONT_CN)
    add_text(slide, "共同盲区：aggregate accuracy 很难定位记忆到底在形成、检索、使用还是更新环节失败。", 0.72, 6.26, 11.92, 0.36, size=17, color=RED, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    slide = add_body_slide(
        prs,
        "StratMem-Bench：从“记住事实”到“战略使用”",
        page,
        "Wu et al., ACL 2026 Long Paper, 657 instances; transformed from LoCoMo with LLM annotation and human validation.",
    )
    add_rect(slide, 0.52, 1.20, 4.08, 4.83, fill=LIGHT, line=LINE, width=1.0)
    add_text(slide, "Input", 0.82, 1.48, 0.86, 0.30, size=17, color=NAVY, bold=True, font=FONT_EN)
    add_text(slide, "xᵢ = (hᵢ, qᵢ, Pᵢ, Mᵢ)", 1.62, 1.44, 2.45, 0.40, size=19, color=INK, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    memory_types = [
        ("must", "保证事实正确，缺失会导致错误", BLUE, BLUE_LIGHT),
        ("nice", "非必需但可增加个性、共情与连贯", GREEN, GREEN_LIGHT),
        ("irr", "与当前目标无关，应主动抑制", RED, RED_LIGHT),
    ]
    for idx, (tag, desc, accent, fill) in enumerate(memory_types):
        y = 2.20 + idx * 0.96
        add_tag(slide, tag, 0.88, y, 0.82, accent=accent, fill=fill)
        add_text(slide, desc, 1.88, y + 0.02, 2.30, 0.42, size=13.2, color=INK, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Output", 0.82, 5.23, 0.90, 0.30, size=17, color=NAVY, bold=True, font=FONT_EN)
    add_text(slide, "角色一致且有策略的回复 ŷ", 1.74, 5.17, 2.42, 0.40, size=14.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 4.92, 1.20, 4.04, 4.83, fill=WHITE, line=NAVY, width=1.1)
    add_text(slide, "四个指标", 5.26, 1.49, 3.36, 0.32, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullet_list(slide, ["SMC：required memory 是否全部满足", "MIQ：使用后的整合质量", "PES：supportive memory 的增益", "CIR：无条件引入 irrelevant memory 的比例"], 5.24, 2.06, 3.36, 2.82, size=14.5)
    add_rect(slide, 9.28, 1.20, 3.55, 4.83, fill=GREEN_LIGHT, line=GREEN, width=1.1)
    add_text(slide, "关键发现", 9.58, 1.49, 2.95, 0.32, size=19, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "模型通常能区分\nmust 与 irr", 9.66, 2.16, 2.78, 0.80, size=20, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "但加入 nice 后，\n“用多少、怎么用”显著变难", 9.66, 3.25, 2.78, 1.02, size=18, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "启示：memory evaluation 要测 selective deployment，不只是 recall。", 9.60, 4.68, 2.90, 0.70, size=14, color=MUTED, align=PP_ALIGN.CENTER)
    page += 1

    slide = add_body_slide(
        prs,
        "近邻工作把评测推向行动，但仍彼此割裂",
        page,
        "Concurrent frontier: Mem2ActBench (ACL 2026), MemFail/AuthMem/SafeCommit (arXiv 2026), MobileMem (technical report, 2026).",
    )
    works = [
        ("Mem2ActBench", "从长期记忆恢复 tool 与参数", "离线 tool call；未纳入 verify / clarify / repair", BLUE, BLUE_LIGHT),
        ("MemFail", "分解 summarization / storage / retrieval", "诊断组件失败；未测行动时的memory regulation", PURPLE, PURPLE_LIGHT),
        ("AuthMem-Bench", "测试 consolidation 后 authority collapse", "聚焦来源权限；未统一覆盖read/use/repair", RED, RED_LIGHT),
        ("SafeCommit", "commit / probe / fallback 下的安全证书", "proof-of-concept latent worlds；未接真实memory stack", GREEN, GREEN_LIGHT),
        ("MobileMem", "一年尺度移动经历与多模态个人记忆", "覆盖更新/偏好，但不是受控matched intervention", ORANGE, RGBColor(255, 245, 220)),
    ]
    y = 1.18
    for name, strength, limitation, accent, fill in works:
        add_rect(slide, 0.58, y, 12.12, 0.91, fill=fill, line=accent, width=1.0)
        add_text(slide, name, 0.82, y + 0.23, 2.05, 0.30, size=15, color=accent, bold=True, font=FONT_EN)
        add_text(slide, strength, 3.02, y + 0.17, 3.44, 0.44, size=13.3, color=INK, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, limitation, 6.67, y + 0.14, 5.60, 0.50, size=12.7, color=MUTED, valign=MSO_ANCHOR.MIDDLE)
        y += 1.03
    add_text(slide, "机会不是再拼一个更大的任务集合，而是用 matched memory-state family 连接 diagnosis → control → repair。", 0.76, 6.35, 11.80, 0.36, size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    add_section_slide(prs, "MemReadyBench")

    slide = add_body_slide(
        prs,
        "核心问题：记忆状态何时足以支持下一步行动？",
        page,
        "Canonical Stage-I proposal: docs/ideas/2026-09-03-memreadybench-unified-research-proposal.md.",
    )
    add_rect(slide, 0.68, 1.18, 11.98, 0.94, fill=BLUE_LIGHT, line=NAVY, width=1.1)
    add_text(slide, "在持久记忆可能陈旧、冲突、不完整或受污染时，Agent 能否诊断其可用性、调控其影响，并在验证后修复记忆以支持未来会话？", 0.95, 1.38, 11.45, 0.46, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    rqs = [
        ("RQ1", "Memory-state diagnosis", "Agent 能否识别当前记忆是 verified、stale、conflicting、incomplete 还是 corrupted？", BLUE, BLUE_LIGHT),
        ("RQ2", "Memory-use control", "在不同风险与证据状态下，能否选择 USE / SUPPRESS / CLARIFY / VERIFY / ABSTAIN？", GREEN, GREEN_LIGHT),
        ("RQ3", "Prospective repair", "获得权威反馈后，能否修复持久记忆，并在下一session稳定复用？", PURPLE, PURPLE_LIGHT),
    ]
    for idx, (rq, title, body, accent, fill) in enumerate(rqs):
        x = 0.68 + idx * 4.08
        add_rect(slide, x, 2.62, 3.72, 2.68, fill=fill, line=accent, width=1.2)
        add_tag(slide, rq, x + 0.30, 2.88, 0.80, accent=accent, fill=WHITE)
        add_text(slide, title, x + 0.26, 3.47, 3.20, 0.40, size=17, color=accent, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(slide, body, x + 0.30, 4.07, 3.12, 0.88, size=14, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, "与StratMem的边界：它测“给定memory池如何表达”；这里测“memory状态变化是否改变工具决策，并能否被修复”。", 0.88, 5.87, 11.52, 0.54, size=16.5, color=RED, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    slide = add_body_slide(
        prs,
        "Stage-I Benchmark总体框架",
        page,
        "Framework source: docs/figures/2026-09-03-memreadybench-stage1-benchmark-framework.png.",
    )
    slide.shapes.add_picture(str(FRAMEWORK), Inches(0.47), Inches(1.05), width=Inches(12.40), height=Inches(6.20))
    page += 1

    slide = add_body_slide(
        prs,
        "Benchmark输入、过程与输出",
        page,
        "Stage-I contains evaluation tasks only; no training/RL module is included.",
    )
    io_steps = [
        ("1", "Canonical world", "构造真实世界 w*、历史 H、用户目标 q 与允许动作集合 A。", BLUE, BLUE_LIGHT),
        ("2", "Matched memory states", "固定任务与世界，仅干预 M：verified / stale / conflict / incomplete / corrupted。", PURPLE, PURPLE_LIGHT),
        ("3", "Black-box agent", "接入被测 memory backend + LLM agent；不泄露memory标签。", TEAL, TEAL_LIGHT),
        ("4", "Action trajectory", "记录 retrieve、use/suppress、clarify/verify、commit/abstain 与tool结果。", GREEN, GREEN_LIGHT),
        ("5", "Session-C probe", "权威反馈后进入新会话，检查修复后的memory能否稳定影响未来行动。", ORANGE, RGBColor(255, 245, 220)),
    ]
    y = 1.14
    for num, title, body, accent, fill in io_steps:
        add_rect(slide, 0.66, y, 12.00, 0.92, fill=fill, line=accent, width=1.0)
        add_text(slide, num, 0.86, y + 0.20, 0.46, 0.42, size=21, color=accent, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(slide, title, 1.50, y + 0.19, 2.25, 0.34, size=15, color=accent, bold=True, font=FONT_EN)
        add_text(slide, body, 3.92, y + 0.13, 8.25, 0.50, size=14, color=INK, valign=MSO_ANCHOR.MIDDLE)
        y += 1.04
    add_text(slide, "核心控制变量：同一实例族只改变 memory state，因此能把行为变化归因于记忆，而不是任务难度。", 0.83, 6.38, 11.66, 0.35, size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    slide = add_body_slide(
        prs,
        "四类任务：同一实例族中的诊断、控制与修复",
        page,
        "Task design follows the benchmark-only Stage-I specification.",
    )
    tasks = [
        ("T1", "Memory Diagnosis", "判断每条memory的状态、来源与缺失证据；输出结构化标签。", BLUE, BLUE_LIGHT),
        ("T2", "Memory-Use Control", "在matched states下选择 USE / SUPPRESS / CLARIFY / VERIFY / ABSTAIN。", GREEN, GREEN_LIGHT),
        ("T3", "Closure & Execution", "证据充分时才提交带副作用的动作；否则执行低风险信息获取。", RED, RED_LIGHT),
        ("T4", "Memory Repair", "验证世界后更新/删除/降权旧memory，并在Session C测试未来复用。", PURPLE, PURPLE_LIGHT),
    ]
    coords = [(0.62, 1.25), (6.78, 1.25), (0.62, 3.80), (6.78, 3.80)]
    for (tag, title, body, accent, fill), (x, y) in zip(tasks, coords):
        add_rect(slide, x, y, 5.90, 2.10, fill=fill, line=accent, width=1.2)
        add_tag(slide, tag, x + 0.28, y + 0.24, 0.76, accent=accent, fill=WHITE)
        add_text(slide, title, x + 1.20, y + 0.23, 4.25, 0.36, size=18, color=accent, bold=True, font=FONT_EN)
        add_text(slide, body, x + 0.34, y + 0.94, 5.22, 0.76, size=15, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, "四类任务不是四个互不相关的数据集，而是同一memory episode上的逐层探针。", 0.96, 6.25, 11.36, 0.38, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    slide = add_body_slide(
        prs,
        "五个主指标：同时看状态、行为与未来收益",
        page,
        "Metric definitions are pre-registered in the Stage-I benchmark specification.",
    )
    metrics = [
        ("MSFA", "Memory-State Factor Accuracy", "对memory状态因子的宏平均准确率；避免多数类掩盖stale/conflict失败。", BLUE, BLUE_LIGHT),
        ("VCS", "Verified Closure Score", "只有在证据充分且行动正确时得分；惩罚premature commit。", GREEN, GREEN_LIGHT),
        ("PMCR", "Paired Memory Causal Responsiveness", "同一实例不同memory state下，策略是否按预期改变。", PURPLE, PURPLE_LIGHT),
        ("LRU", "Lifecycle Repair Utility", "修复前后在Session C的未来任务收益差。", ORANGE, RGBColor(255, 245, 220)),
        ("NAR", "Negative-memory Avoidance Rate", "面对错误/污染记忆时，成功抑制其影响的比例。", RED, RED_LIGHT),
    ]
    y = 1.13
    for acronym, name, body, accent, fill in metrics:
        add_rect(slide, 0.62, y, 12.08, 0.94, fill=fill, line=accent, width=1.0)
        add_text(slide, acronym, 0.88, y + 0.18, 1.05, 0.38, size=19, color=accent, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(slide, name, 2.12, y + 0.20, 3.20, 0.34, size=14.5, color=accent, bold=True, font=FONT_EN)
        add_text(slide, body, 5.48, y + 0.13, 6.76, 0.50, size=13.5, color=INK, valign=MSO_ANCHOR.MIDDLE)
        y += 1.04
    add_text(slide, "主张成立的最低条件：不仅总体分下降，还要在 matched pairs 上出现方向正确、统计稳定的策略变化。", 0.82, 6.38, 11.68, 0.36, size=16.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    slide = add_body_slide(
        prs,
        "最小可行实验与Go / No-Go",
        page,
        "Pilot goal: validate the benchmark's diagnostic signal before scaling data or proposing a training method.",
    )
    add_rect(slide, 0.64, 1.20, 6.00, 4.90, fill=LIGHT, line=LINE, width=1.0)
    add_text(slide, "Pilot配置", 0.98, 1.49, 5.30, 0.36, size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullet_list(
        slide,
        [
            "200-300个canonical episodes，每个生成3-5种matched memory states。",
            "3类领域：日程/沟通、文件与知识工作、受约束的工具调用。",
            "Baselines：long-context、RAG、直接memory reader、StratMem-style selector、SafeCommit wrapper。",
            "至少3个backbone + 2个memory backend；报告效果、成本、延迟与拒答率。",
        ],
        0.98,
        2.10,
        5.24,
        3.37,
        size=14.2,
    )
    add_rect(slide, 6.96, 1.20, 5.72, 2.18, fill=GREEN_LIGHT, line=GREEN, width=1.2)
    add_text(slide, "GO", 7.27, 1.50, 0.90, 0.45, size=26, color=GREEN, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(slide, "memory-state intervention 产生显著、可复现的PMCR；错误能被归因到diagnosis/control/repair中的具体环节。", 8.25, 1.46, 3.96, 1.05, size=14.5, color=INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 6.96, 3.92, 5.72, 2.18, fill=RED_LIGHT, line=RED, width=1.2)
    add_text(slide, "NO-GO", 7.16, 4.24, 1.20, 0.45, size=23, color=RED, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(slide, "模型对memory state不敏感、指标主要复刻最终准确率，或自动标注无法被人工稳定复核；此时先收缩问题而不是扩数据。", 8.35, 4.11, 3.82, 1.18, size=14.2, color=INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Stage II训练暂不进入框架图：只有Stage I证明存在稳定诊断信号后，才设计memory controller或RL。", 0.80, 6.37, 11.72, 0.35, size=16.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    page += 1

    slide = add_body_slide(prs, "参考文献：LLM-based KGR", page)
    kgr_left = [
        "Pan, S. et al. (2024). Unifying Large Language Models and Knowledge Graphs: A Roadmap. IEEE TKDE. doi:10.1109/TKDE.2024.3352100.",
        "Sun, J. et al. (2024). Think-on-Graph: Deep and Responsible Reasoning of Large Language Model on Knowledge Graph. ICLR 2024. openreview.net/forum?id=nnVO1PvbTv.",
        "Luo, L., Li, Y.-F., Haffari, G., & Pan, S. (2024). Reasoning on Graphs: Faithful and Interpretable Large Language Model Reasoning. ICLR 2024. openreview.net/forum?id=ZGNWW7xZ6Q.",
        "Jiang, J. et al. (2025). KG-Agent: An Efficient Autonomous Agent Framework for Complex Reasoning over Knowledge Graph. ACL 2025, 9505-9523. doi:10.18653/v1/2025.acl-long.468.",
    ]
    kgr_right = [
        "Wu, J., Zhu, J., Liu, Y., Xu, M., & Jin, Y. (2025). Agentic Reasoning: A Streamlined Framework for Enhancing LLM Reasoning with Agentic Tools. ACL 2025, 28489-28503. doi:10.18653/v1/2025.acl-long.1383.",
        "Liu, Y. et al. (2026). PathMind: A Retrieve-Prioritize-Reason Framework for Knowledge Graph Reasoning with Large Language Models. AAAI 40(18), 15386-15393. doi:10.1609/aaai.v40i18.38565.",
        "Yan, S. et al. (2026). Explore-on-Graph: Incentivizing Autonomous Exploration of LLMs on KGs with Path-refined Reward Modeling. ICLR 2026. arXiv:2602.21728.",
        "Zhang, Y. et al. (2026). Backjump-on-Graph: Empowering LLMs with Reinforced Retrospective Exploration for Agentic KGR. ICML 2026, PMLR 306.",
    ]
    add_reference_columns(slide, kgr_left, kgr_right)
    add_citation(slide, "All venue/status metadata verified against official proceedings, ACL Anthology, OpenReview, or arXiv on 2026-09-03.", width=9.8)
    page += 1

    slide = add_body_slide(prs, "参考文献：Agent Memory Benchmarks", page)
    memory_left = [
        "Yang, X. et al. (2024). CRAG - Comprehensive RAG Benchmark. NeurIPS 2024 Datasets and Benchmarks Track, 10470-10490. doi:10.52202/079017-0335.",
        "Maharana, A. et al. (2024). Evaluating Very Long-Term Conversational Memory of LLM Agents. ACL 2024, 13851-13870. doi:10.18653/v1/2024.acl-long.747.",
        "Wu, D. et al. (2025). LongMemEval: Benchmarking Chat Assistants on Long-Term Interactive Memory. ICLR 2025. arXiv:2410.10813.",
        "He, J. et al. (2025). MADial-Bench: Towards Real-World Evaluation of Memory-Augmented Dialogue Generation. NAACL 2025, 9902-9921.",
        "Hu, Y. et al. (2025). Memory in the Age of AI Agents. arXiv:2512.13564.",
    ]
    memory_right = [
        "Wu, Y., Wu, T., Zhu, M., Sha, H., & Wang, H. (2026). StratMem-Bench: Evaluating Strategic Memory Use in Virtual Character Conversation Beyond Factual Recall. ACL 2026, 32309-32328. doi:10.18653/v1/2026.acl-long.1491.",
        "Shen, Y., Li, K., Zhou, W., & Hu, S. (2026). Mem2ActBench: Evaluating Long-Term Memory Utilization in Task-Oriented Autonomous Agents. ACL 2026. aclanthology.org/2026.acl-long.370.",
        "Garg, I., Kolhe, N., Song, D., & Zhao, X. (2026). MemFail: Stress-Testing Failure Modes of LLM Memory Systems. arXiv:2605.26667.",
        "Zhan, Q. et al. (2026). When Memory Becomes Authority: Benchmarking Authority Collapse at the Memory Consolidation Boundary. arXiv:2608.01679.",
        "Akewar, M., & Ranjan, R. (2026). SafeCommit: Certifying When Memory-Grounded Agents May Safely Act. arXiv:2608.04289.",
    ]
    add_reference_columns(slide, memory_left, memory_right)
    add_citation(slide, "Classic-work selection follows the StratMem-Bench introduction/related work; concurrent references are labeled as arXiv where applicable.", width=10.5)

    # Keep the original closing slide as the final slide.
    prs.slides._sldIdLst.remove(closing_sld_id)
    prs.slides._sldIdLst.append(closing_sld_id)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
