# -*- coding: utf-8 -*-
"""Build the MetaMemBench research proposal deck.

Run from the repository root:
    python scripts/build_metamembench_pptx.py

The script uses python-pptx. For the local workspace, dependencies are installed
under .tmp/pptx_deps; a normal environment can simply install python-pptx.
"""

from __future__ import annotations

import sys
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
LOCAL_DEPS = ROOT / ".tmp" / "pptx_deps"
if LOCAL_DEPS.exists():
    sys.path.insert(0, str(LOCAL_DEPS))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = ROOT / "docs" / "slides" / "2026-09-03-metamembench-research-proposal.pptx"

SW = 13.333
SH = 7.5
FONT = "Microsoft YaHei"
FONT_EN = "Aptos"

INK = "172033"
MUTED = "5B667A"
LIGHT = "F5F7FA"
WHITE = "FFFFFF"
LINE = "D9E0E8"
TEAL = "0F766E"
TEAL_LIGHT = "DDF3EF"
BLUE = "2563EB"
BLUE_LIGHT = "E7EEFF"
ORANGE = "D97706"
ORANGE_LIGHT = "FFF0D5"
RED = "B42318"
RED_LIGHT = "FDE8E7"
GREEN = "16803A"
GREEN_LIGHT = "E4F4E8"
PURPLE = "6D4BC3"
PURPLE_LIGHT = "EFE9FF"
GRAY_LIGHT = "EDF1F5"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=False, width=1):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    return shape


def line(slide, x, y, w, h=0.02, color=LINE):
    return rect(slide, x, y, w, h, fill=color, line=color)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font=FONT,
    margin=0.03,
    line_spacing=1.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    parts = str(text).split("\n")
    for i, part in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = part
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, size=18, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.05
    for spec in runs:
        run = p.add_run()
        run.text = spec[0]
        run.font.name = spec[4] if len(spec) > 4 else FONT
        run.font.size = Pt(spec[3] if len(spec) > 3 else size)
        run.font.bold = spec[1]
        run.font.color.rgb = rgb(spec[2])
    return box


def add_list(slide, items, x, y, w, h, size=16, color=INK, bullet_color=TEAL, gap=5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.05
        p.space_after = Pt(gap)
        bullet = p.add_run()
        bullet.text = "• "
        bullet.font.name = FONT
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = rgb(bullet_color)
        body = p.add_run()
        body.text = item
        body.font.name = FONT
        body.font.size = Pt(size)
        body.font.color.rgb = rgb(color)
    return box


def add_badge(slide, text, x, y, w, fill=TEAL_LIGHT, color=TEAL, size=11):
    rect(slide, x, y, w, 0.34, fill=fill, line=fill, radius=True)
    add_text(
        slide,
        text,
        x,
        y + 0.01,
        w,
        0.28,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_card(slide, x, y, w, h, title, body, accent=TEAL, fill=WHITE, title_size=16, body_size=13):
    rect(slide, x, y, w, h, fill=fill, line=LINE, radius=True)
    rect(slide, x, y, 0.07, h, fill=accent, line=accent, radius=False)
    add_text(slide, title, x + 0.20, y + 0.14, w - 0.34, 0.34, size=title_size, color=accent, bold=True)
    add_text(slide, body, x + 0.20, y + 0.55, w - 0.34, h - 0.66, size=body_size, color=INK, line_spacing=1.08)


def add_number_card(slide, x, y, w, h, number, label, accent=TEAL, sub=None):
    rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True)
    add_text(slide, number, x + 0.14, y + 0.13, w - 0.28, 0.44, size=24, color=accent, bold=True)
    add_text(slide, label, x + 0.14, y + 0.59, w - 0.28, 0.30, size=11, color=MUTED, bold=True)
    if sub:
        add_text(slide, sub, x + 0.14, y + 0.91, w - 0.28, h - 1.00, size=9, color=MUTED)


def add_chevron(slide, x, y, w=0.28, h=0.36, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(color)
    return shape


def add_header(slide, section, title, number, subtitle=None):
    rect(slide, 0, 0, SW, SH, fill=LIGHT, line=LIGHT)
    rect(slide, 0, 0, 0.14, SH, fill=TEAL, line=TEAL)
    add_text(slide, section.upper(), 0.54, 0.30, 4.0, 0.28, size=10, color=TEAL, bold=True, font=FONT_EN)
    add_text(slide, title, 0.54, 0.64, 12.1, 0.52, size=25, color=INK, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.56, 1.16, 11.9, 0.36, size=11, color=MUTED)
    line(slide, 0.54, 1.53, 12.22, 0.015, color=LINE)
    add_text(slide, f"{number:02d}", 12.20, 0.28, 0.52, 0.28, size=10, color=MUTED, bold=True, align=PP_ALIGN.RIGHT, font=FONT_EN)


def add_footer(slide, citation=None):
    if citation:
        add_text(slide, citation, 0.56, 7.14, 11.7, 0.20, size=8, color=MUTED, font=FONT_EN)
    add_text(slide, "MetaMemBench · 2026-09-03", 10.90, 7.14, 1.84, 0.20, size=8, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT_EN)


def add_table(slide, x, y, widths, row_h, headers, rows, header_fill=INK, body_size=11):
    cur_x = x
    for j, width in enumerate(widths):
        rect(slide, cur_x, y, width, row_h, fill=header_fill, line=WHITE)
        add_text(slide, headers[j], cur_x + 0.06, y + 0.06, width - 0.12, row_h - 0.10, size=11, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
        cur_x += width
    for i, row in enumerate(rows):
        cur_x = x
        fill = WHITE if i % 2 == 0 else GRAY_LIGHT
        yy = y + row_h * (i + 1)
        for j, width in enumerate(widths):
            rect(slide, cur_x, yy, width, row_h, fill=fill, line=WHITE)
            add_text(slide, row[j], cur_x + 0.06, yy + 0.05, width - 0.12, row_h - 0.08, size=body_size, color=INK, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
            cur_x += width


def add_reference(slide, index, title, venue, url, x, y, w):
    add_text(slide, f"{index:02d}", x, y, 0.35, 0.28, size=10, color=TEAL, bold=True, font=FONT_EN)
    add_text(slide, title, x + 0.42, y, w - 0.42, 0.25, size=10, color=INK, bold=True, font=FONT_EN)
    box = add_text(slide, f"{venue}  ·  {url}", x + 0.42, y + 0.24, w - 0.42, 0.24, size=8, color=MUTED, font=FONT_EN)
    for paragraph in box.text_frame.paragraphs:
        for run in paragraph.runs:
            if "http" in run.text:
                run.hyperlink.address = url


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]

    # 01 Cover
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SW, SH, fill=LIGHT, line=LIGHT)
    rect(slide, 0, 0, SW, 0.16, fill=TEAL, line=TEAL)
    add_badge(slide, "RESEARCH PROPOSAL", 0.72, 0.64, 1.72, fill=TEAL_LIGHT, color=TEAL, size=10)
    add_text(slide, "MetaMemBench", 0.72, 1.20, 7.25, 0.70, size=34, color=INK, bold=True, font=FONT_EN)
    add_text(slide, "面向 Agentic Memory 的\n记忆监控与控制策略诊断基准", 0.72, 1.95, 7.35, 1.42, size=27, color=INK, bold=True, line_spacing=1.0)
    add_text(slide, "Evaluating Monitoring and Control of Agent Memory under Uncertainty", 0.76, 3.51, 7.35, 0.44, size=14, color=TEAL, bold=True, font=FONT_EN)
    rect(slide, 0.72, 4.28, 7.45, 1.40, fill=WHITE, line=LINE, radius=True)
    add_text(slide, "核心问题", 0.98, 4.54, 1.00, 0.30, size=12, color=ORANGE, bold=True)
    add_text(slide, "当前记忆是否足以支撑行动？若不足，Agent 应该继续搜索、向用户澄清，还是拒绝承诺？", 0.98, 4.92, 6.86, 0.60, size=17, color=INK, bold=True, line_spacing=1.08)
    rect(slide, 8.64, 0.76, 4.02, 5.87, fill=INK, line=INK, radius=True)
    add_text(slide, "从“检索到什么”\n走向“何时可以行动”", 9.02, 1.10, 3.25, 0.95, size=22, color=WHITE, bold=True, line_spacing=1.0)
    keywords = [
        ("01", "Sufficiency", "充分性判断", TEAL_LIGHT, TEAL),
        ("02", "Control", "搜 / 问 / 做 / 拒绝", BLUE_LIGHT, BLUE),
        ("03", "Counterfactual", "反事实动作翻转", ORANGE_LIGHT, ORANGE),
        ("04", "Diagnosis", "分解失败来源", RED_LIGHT, RED),
    ]
    for i, (num, en, zh, fill, accent) in enumerate(keywords):
        yy = 2.34 + i * 0.88
        rect(slide, 9.00, yy, 3.30, 0.66, fill=fill, line=fill, radius=True)
        add_text(slide, num, 9.14, yy + 0.16, 0.34, 0.24, size=10, color=accent, bold=True, font=FONT_EN)
        add_text(slide, en, 9.55, yy + 0.10, 1.35, 0.25, size=12, color=accent, bold=True, font=FONT_EN)
        add_text(slide, zh, 10.90, yy + 0.10, 1.22, 0.28, size=10, color=INK, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, "方案汇报 · Agent Memory Benchmark", 0.76, 6.80, 4.50, 0.24, size=10, color=MUTED)
    add_text(slide, "2026-09-03", 11.15, 6.80, 1.14, 0.24, size=10, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT_EN)

    # 02 Executive decision
    slide = prs.slides.add_slide(blank)
    add_header(slide, "01 / Decision", "先给结论：不要只做“动态 StratMem”", 2)
    add_text(slide, "推荐主线", 0.62, 1.82, 2.0, 0.38, size=14, color=TEAL, bold=True)
    rect(slide, 0.60, 2.19, 7.38, 2.20, fill=INK, line=INK, radius=True)
    add_text(slide, "MetaMemBench", 0.92, 2.52, 2.96, 0.42, size=25, color=WHITE, bold=True, font=FONT_EN)
    add_text(slide, "评测 Agent 是否知道当前记忆“够不够”，并在承诺行动前选择正确的控制动作。", 0.92, 3.09, 6.50, 0.82, size=20, color=WHITE, bold=True, line_spacing=1.08)
    add_badge(slide, "diagnostic", 0.92, 4.62, 1.24, fill=TEAL_LIGHT, color=TEAL)
    add_badge(slide, "counterfactual", 2.28, 4.62, 1.52, fill=BLUE_LIGHT, color=BLUE)
    add_badge(slide, "sufficiency-aware", 3.92, 4.62, 1.78, fill=ORANGE_LIGHT, color=ORANGE)
    add_badge(slide, "control-before-commitment", 5.82, 4.62, 2.16, fill=RED_LIGHT, color=RED, size=9)
    add_text(slide, "暂时放弃的表述", 8.46, 1.82, 2.2, 0.38, size=14, color=RED, bold=True)
    add_card(slide, 8.44, 2.19, 4.26, 1.35, "Dynamic Strategic Memory", "多轮化本身容易被视为 StratMem-Bench 的自然扩展。", accent=RED, fill=RED_LIGHT, title_size=15, body_size=12)
    add_card(slide, 8.44, 3.75, 4.26, 1.35, "“首个 memory control”", "AgeMem、MemCon、MemSearcher 已学习何时与如何调用 memory。", accent=ORANGE, fill=ORANGE_LIGHT, title_size=15, body_size=12)
    add_card(slide, 8.44, 5.31, 4.26, 1.25, "真正可争取的贡献", "统一、可诊断、可执行的 memory-control 评测协议。", accent=TEAL, fill=TEAL_LIGHT, title_size=15, body_size=12)
    add_footer(slide, "Positioning based on Memory in the Age of AI Agents; StratMem-Bench; MemCon; AgeMem.")

    # 03 Positioning
    slide = prs.slides.add_slide(blank)
    add_header(slide, "02 / Positioning", "它属于 Agentic Memory，但与 RAG / Context Engineering 交叉", 3)
    areas = [
        (0.70, 2.00, 3.00, 1.28, "Agent Memory", "持久、跨 session、随交互演化", TEAL, TEAL_LIGHT),
        (0.70, 4.02, 3.00, 1.28, "Agentic RAG", "主动检索、查询改写、停止策略", BLUE, BLUE_LIGHT),
        (9.63, 2.00, 3.00, 1.28, "LLM Memory", "参数 / KV / 长上下文机制", PURPLE, PURPLE_LIGHT),
        (9.63, 4.02, 3.00, 1.28, "Context Engineering", "工具、预算与临时上下文组织", ORANGE, ORANGE_LIGHT),
    ]
    for x, y, w, h, title, body, accent, fill in areas:
        add_card(slide, x, y, w, h, title, body, accent=accent, fill=fill, title_size=16, body_size=11)
    rect(slide, 4.15, 2.02, 4.98, 3.80, fill=WHITE, line=TEAL, radius=True, width=2)
    add_badge(slide, "OUR POSITION", 5.68, 2.30, 1.90, fill=TEAL_LIGHT, color=TEAL, size=10)
    add_text(slide, "Persistent memory", 5.05, 2.92, 3.18, 0.33, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_chevron(slide, 6.50, 3.34, 0.28, 0.30, TEAL)
    add_text(slide, "Monitor sufficiency", 5.05, 3.70, 3.18, 0.33, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_chevron(slide, 6.50, 4.13, 0.28, 0.30, TEAL)
    add_text(slide, "Control next action", 5.05, 4.49, 3.18, 0.33, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_chevron(slide, 6.50, 4.92, 0.28, 0.30, TEAL)
    add_text(slide, "Verify downstream consequence", 4.72, 5.29, 3.84, 0.33, size=14, color=TEAL, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(slide, "边界条件：必须同时保留 persistent、interaction-derived memory，Agent-controlled access，以及可验证的行动后果。", 0.72, 6.25, 11.95, 0.48, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Hu et al., Memory in the Age of AI Agents, arXiv:2512.13564.")

    # 04 Evolution
    slide = prs.slides.add_slide(blank)
    add_header(slide, "03 / Prior Work", "已有工作如何一步步逼近“记忆控制”", 4)
    line(slide, 1.00, 3.47, 11.30, 0.05, color=LINE)
    stages = [
        (0.70, "1", "Recall", "能否找回事实", "LoCoMo\nLongMemEval", BLUE, BLUE_LIGHT),
        (3.00, "2", "Strategic Use", "该用哪些记忆", "StratMem-Bench", TEAL, TEAL_LIGHT),
        (5.30, "3", "Memory → Action", "能否驱动工具/任务", "Mem2ActBench\nMemoryArena", ORANGE, ORANGE_LIGHT),
        (7.60, "4", "Lifecycle / Control", "何时检索、更新、忘记", "MemOps\nAgeMem · MemCon", PURPLE, PURPLE_LIGHT),
        (9.90, "5", "Control Diagnosis", "记忆够不够，下一步做什么", "MetaMemBench", RED, RED_LIGHT),
    ]
    for x, num, title, question, refs, accent, fill in stages:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.72), Inches(3.18), Inches(0.58), Inches(0.58))
        circ.fill.solid(); circ.fill.fore_color.rgb = rgb(accent); circ.line.color.rgb = rgb(accent)
        add_text(slide, num, x + 0.72, 3.27, 0.58, 0.24, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
        rect(slide, x, 1.90, 2.02, 1.02, fill=fill, line=accent, radius=True)
        add_text(slide, title, x + 0.10, 2.08, 1.82, 0.28, size=14, color=accent, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(slide, question, x + 0.10, 2.42, 1.82, 0.30, size=11, color=INK, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, refs, x + 0.05, 4.04, 1.92, 0.82, size=11, color=MUTED, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    rect(slide, 0.70, 5.42, 11.85, 1.06, fill=INK, line=INK, radius=True)
    add_text(slide, "现状：方法已经在学习 memory action；评测却仍主要停留在候选池选择或最终任务成功。", 1.02, 5.72, 11.15, 0.40, size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Representative works: StratMem-Bench; Mem2ActBench; MemoryArena; MemOps; AgeMem; MemCon.")

    # 05 StratMem deep dive
    slide = prs.slides.add_slide(blank)
    add_header(slide, "03 / Prior Work", "StratMem-Bench：重要起点，但它问的是“用哪些”", 5)
    add_number_card(slide, 0.65, 1.84, 1.75, 1.20, "657", "INSTANCES", TEAL)
    add_number_card(slide, 2.55, 1.84, 1.75, 1.20, "8.79", "AVG. POOL", BLUE)
    add_number_card(slide, 4.45, 1.84, 1.75, 1.20, "< 50%", "MUST+NICE SMC", RED)
    add_text(slide, "输入", 0.72, 3.40, 0.62, 0.30, size=12, color=MUTED, bold=True)
    rect(slide, 1.22, 3.28, 1.42, 0.58, fill=GRAY_LIGHT, line=LINE, radius=True)
    add_text(slide, "固定候选池", 1.30, 3.43, 1.26, 0.24, size=12, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_chevron(slide, 2.78, 3.38, 0.30, 0.35, TEAL)
    for i, (name, col, fill) in enumerate([("must", RED, RED_LIGHT), ("nice", TEAL, TEAL_LIGHT), ("irrelevant", MUTED, GRAY_LIGHT)]):
        rect(slide, 3.22 + i * 0.98, 3.28, 0.84, 0.58, fill=fill, line=col, radius=True)
        add_text(slide, name, 3.26 + i * 0.98, 3.43, 0.76, 0.24, size=9, color=col, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_card(slide, 0.66, 4.32, 5.77, 1.76, "它已经解决", "query-conditioned 记忆角色不是固定属性；模型需要使用必要记忆、选择性利用支持记忆，并抑制无关记忆。", accent=TEAL, fill=TEAL_LIGHT, title_size=17, body_size=14)
    add_card(slide, 6.68, 1.84, 6.02, 4.24, "它没有要求 Agent 做什么", "• 判断当前 packet 是否足以安全行动\n• 主动搜索缺失记忆或改写 query\n• 区分“store 里有”与“只能问用户”\n• 在 EXECUTE / ASK / ABSTAIN 间控制承诺\n• 让动作进入环境并验证后果", accent=RED, fill=WHITE, title_size=18, body_size=15)
    rect(slide, 0.66, 6.26, 12.04, 0.52, fill=INK, line=INK, radius=True)
    add_text(slide, "StratMem asks which memories should appear; MetaMemBench asks whether memory can justify acting at all.", 0.94, 6.38, 11.50, 0.24, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_footer(slide, "Wu et al., StratMem-Bench, ACL 2026 Main Conference, Long Paper.")

    # 06 Nearest work matrix
    slide = prs.slides.add_slide(blank)
    add_header(slide, "03 / Prior Work", "最邻近工作：各自解决了一块，但没有统一诊断控制决策", 6)
    rows = [
        ("Mem2ActBench", "长记忆 → 工具参数", "离线 tool-call；搜索/澄清不是正式动作"),
        ("MemoryArena", "跨 session 环境闭环", "只有 task success/progress，缺 decision-point gold"),
        ("AMemGym", "on-policy 长程对话", "动态环境已有；memory-control 不是独立评测单元"),
        ("MemOps", "remember/update/forget/reflect", "评 lifecycle，不评 query-time evidence sufficiency"),
        ("TANGLE / StateMem", "冲突与状态更新", "覆盖特定失效类型，不覆盖完整控制状态空间"),
        ("AgeMem", "工具化 memory operation + RL", "证明可学控制；缺独立统一的动作诊断协议"),
        ("MemCon", "Memory MDP + 在线策略", "依赖终局反馈，难解释 monitor/action 为何对错"),
    ]
    add_table(slide, 0.65, 1.82, [2.18, 3.32, 6.52], 0.61, ["工作", "已经解决", "仍留下的可验证问题"], rows, body_size=11)
    rect(slide, 0.65, 6.74, 12.02, 0.30, fill=TEAL_LIGHT, line=TEAL_LIGHT)
    add_text(slide, "MetaMemBench 不替代这些任务，而是提供跨系统的 monitor / control / outcome / cost 诊断层。", 0.82, 6.77, 11.70, 0.20, size=10, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Sources: ACL 2026; ICLR 2026; arXiv 2026 papers listed in the proposal.")

    # 07 Gap
    slide = prs.slides.add_slide(blank)
    add_header(slide, "04 / Research Gap", "现有评测从两侧逼近，却漏掉了“行动前控制”", 7)
    add_card(slide, 0.70, 1.92, 3.72, 2.20, "固定候选池评测", "看得见全部候选\n能评“选哪些”\n不能评搜索、澄清、停止与拒绝", accent=BLUE, fill=BLUE_LIGHT, title_size=18, body_size=15)
    add_card(slide, 8.92, 1.92, 3.72, 2.20, "端到端任务评测", "有真实行动后果\n但最终成功混合了检索、控制、执行\n甚至可能靠猜测“碰巧成功”", accent=ORANGE, fill=ORANGE_LIGHT, title_size=18, body_size=15)
    add_chevron(slide, 4.63, 2.81, 0.52, 0.52, BLUE)
    add_chevron(slide, 8.18, 2.81, 0.52, 0.52, ORANGE)
    rect(slide, 5.12, 1.88, 3.10, 2.30, fill=INK, line=INK, radius=True)
    add_text(slide, "缺失的评测层", 5.45, 2.20, 2.44, 0.34, size=16, color=TEAL_LIGHT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Memory-state\nmonitoring\n+\nnext-action control", 5.49, 2.72, 2.36, 1.13, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN, line_spacing=0.9)
    add_text(slide, "RESEARCH GAP", 0.74, 4.68, 2.20, 0.28, size=11, color=RED, bold=True, font=FONT_EN)
    rect(slide, 0.70, 5.04, 11.94, 1.38, fill=WHITE, line=RED, radius=True, width=2)
    add_text(slide, "缺少一种受控、可诊断的评测协议，用来测试 Agent 是否能判断当前记忆状态的充分性，并在承诺回答或执行动作之前，选择合适的下一步控制动作。", 1.02, 5.35, 11.30, 0.68, size=19, color=INK, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.05)
    add_footer(slide, "Gap synthesized from StratMem-Bench, Mem2ActBench, MemoryArena, MemTrace, AgeMem and MemCon.")

    # 08 Formal problem
    slide = prs.slides.add_slide(blank)
    add_header(slide, "05 / Problem", "正式问题：从记忆状态 zₜ 映射到可接受控制动作", 8)
    add_text(slide, "六类隐藏状态", 0.68, 1.80, 2.10, 0.34, size=14, color=TEAL, bold=True)
    states = [
        ("NO MEMORY NEEDED", "无需历史即可安全执行", GREEN, GREEN_LIGHT),
        ("SUFFICIENT", "已有最小充分证据闭包", TEAL, TEAL_LIGHT),
        ("RETRIEVABLE MISSING", "缺失证据仍在 store 中", BLUE, BLUE_LIGHT),
        ("USER-ONLY MISSING", "只能由用户补充约束", ORANGE, ORANGE_LIGHT),
        ("STALE / CONFLICT", "旧值或冲突仍可解析", PURPLE, PURPLE_LIGHT),
        ("IRREDUCIBLE", "现有来源无法唯一支持行动", RED, RED_LIGHT),
    ]
    for i, (name, body, accent, fill) in enumerate(states):
        col = i % 2
        row = i // 2
        x = 0.68 + col * 3.18
        y = 2.18 + row * 1.16
        rect(slide, x, y, 2.88, 0.90, fill=fill, line=accent, radius=True)
        add_text(slide, name, x + 0.13, y + 0.12, 2.62, 0.25, size=10, color=accent, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(slide, body, x + 0.13, y + 0.46, 2.62, 0.23, size=10, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "四种外部动作", 7.23, 1.80, 2.10, 0.34, size=14, color=ORANGE, bold=True)
    actions = [
        ("SEARCH_MEMORY", "继续检索 / 改写 query", BLUE, BLUE_LIGHT),
        ("ASK_USER", "询问缺失 slot", ORANGE, ORANGE_LIGHT),
        ("EXECUTE", "回答或调用工具，并提交 evidence IDs", GREEN, GREEN_LIGHT),
        ("ABSTAIN", "拒绝不受支持的承诺", RED, RED_LIGHT),
    ]
    for i, (name, body, accent, fill) in enumerate(actions):
        y = 2.18 + i * 0.91
        add_card(slide, 7.22, y, 5.42, 0.78, name, body, accent=accent, fill=fill, title_size=13, body_size=10)
    rect(slide, 0.68, 5.88, 11.96, 0.84, fill=INK, line=INK, radius=True)
    add_text(slide, "关键方法学约束", 0.96, 6.08, 1.48, 0.26, size=12, color=TEAL_LIGHT, bold=True)
    add_text(slide, "gold 只能依据 Agent 当时可见的信息；不确定时用 admissible action set + cost/regret，而不是强迫唯一动作。", 2.52, 6.06, 9.78, 0.32, size=14, color=WHITE, bold=True)
    add_footer(slide, "Formal state/action taxonomy proposed in MetaMemBench.")

    # 09 Challenges
    slide = prs.slides.add_slide(blank)
    add_header(slide, "06 / Challenges", "真正难点不是把对话变长，而是把“可行动性”标对、测准", 9)
    challenges = [
        (0.68, 1.90, "01", "充分性 ≠ 相关性", "单条都相关，不代表组合已经覆盖全部约束；旧偏好也可能高度相似却有害。", TEAL, TEAL_LIGHT),
        (6.78, 1.90, "02", "部分可观测下的动作 gold", "从当前 packet 未必能知道 store 是否仍有答案；SEARCH 与 ASK 可能都合理。", BLUE, BLUE_LIGHT),
        (0.68, 3.45, "03", "控制失败的因果归因", "最终错误可能来自 monitor、retriever、controller 或 executor，单一成功率无法区分。", ORANGE, ORANGE_LIGHT),
        (6.78, 3.45, "04", "风险与成本同时校准", "过早执行有风险；反复搜索、频繁追问也会伤害效率与用户体验。", RED, RED_LIGHT),
    ]
    for x, y, num, title, body, accent, fill in challenges:
        rect(slide, x, y, 5.86, 1.31, fill=fill, line=accent, radius=True)
        add_text(slide, num, x + 0.20, y + 0.22, 0.48, 0.34, size=15, color=accent, bold=True, font=FONT_EN)
        add_text(slide, title, x + 0.82, y + 0.18, 4.76, 0.31, size=16, color=accent, bold=True)
        add_text(slide, body, x + 0.82, y + 0.58, 4.72, 0.50, size=12, color=INK, line_spacing=1.05)
    rect(slide, 0.68, 5.22, 11.96, 1.13, fill=WHITE, line=PURPLE, radius=True, width=2)
    add_text(slide, "05", 0.92, 5.48, 0.48, 0.34, size=15, color=PURPLE, bold=True, font=FONT_EN)
    add_text(slide, "数据有效性与捷径控制", 1.55, 5.40, 2.86, 0.32, size=16, color=PURPLE, bold=True)
    add_text(slide, "需要 symbolic world、反事实配对、base-task 分割、表述改写与 human challenge split，防止模型只记模板。", 4.48, 5.38, 7.78, 0.48, size=13, color=INK, bold=True)
    add_footer(slide, "Design risks informed by Mem2ActBench, TANGLE, StateMemBench and MemTrace.")

    # 10 Benchmark tracks
    slide = prs.slides.add_slide(blank)
    add_header(slide, "07 / Benchmark", "三个递进 Track：先隔离能力，再开放控制，最后验证行动", 10)
    tracks = [
        (0.66, "A", "Oracle Packet Monitoring", "给定完整可访问证据 + distractors\n不开放搜索\n隔离 monitor / reader", TEAL, TEAL_LIGHT),
        (4.48, "B", "Active Memory Control", "仅给任务 + memory API\n多轮 search / rewrite / stop / ask\n评 timing、intent 与预算", BLUE, BLUE_LIGHT),
        (8.30, "C", "Closed-Loop Memory-to-Action", "动作进入确定性工具环境\n精确验证状态变化\n排除语言 judge 偏差", ORANGE, ORANGE_LIGHT),
    ]
    for x, letter, title, body, accent, fill in tracks:
        rect(slide, x, 1.90, 3.48, 3.18, fill=fill, line=accent, radius=True, width=2)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.20), Inches(2.15), Inches(0.62), Inches(0.62))
        circ.fill.solid(); circ.fill.fore_color.rgb = rgb(accent); circ.line.color.rgb = rgb(accent)
        add_text(slide, letter, x + 0.20, 2.27, 0.62, 0.24, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(slide, title, x + 0.98, 2.13, 2.22, 0.60, size=16, color=accent, bold=True, font=FONT_EN)
        add_text(slide, body, x + 0.24, 3.05, 3.00, 1.55, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.15)
    add_chevron(slide, 4.15, 3.20, 0.28, 0.40, TEAL)
    add_chevron(slide, 7.97, 3.20, 0.28, 0.40, BLUE)
    add_text(slide, "两个核心任务域", 0.68, 5.46, 2.10, 0.30, size=13, color=MUTED, bold=True)
    add_card(slide, 0.68, 5.84, 5.78, 0.90, "Personalized Tool Execution", "日历 / 出行 / 购物 / 提醒 / 文件操作；gold 为结构化 tool call 与环境终态。", accent=GREEN, fill=GREEN_LIGHT, title_size=14, body_size=11)
    add_card(slide, 6.68, 5.84, 5.96, 0.90, "Progressive Search and Planning", "跨 session 约束与反馈；gold 为 evidence closure、规划约束与最终状态。", accent=PURPLE, fill=PURPLE_LIGHT, title_size=14, body_size=11)
    add_footer(slide, "Track design bridges controlled diagnosis and executable agent evaluation.")

    # 11 Data engine
    slide = prs.slides.add_slide(blank)
    add_header(slide, "08 / Data", "数据构造：先建可执行世界，再生成自然语言", 11)
    pipeline = [
        ("1", "Symbolic\nworld", TEAL),
        ("2", "Cross-session\nhistory", BLUE),
        ("3", "Requirement\ngraph", PURPLE),
        ("4", "Counterfactual\nvariants", ORANGE),
        ("5", "Language\nrendering", GREEN),
        ("6", "Deterministic\nvalidation", RED),
    ]
    for i, (num, label, accent) in enumerate(pipeline):
        x = 0.66 + i * 2.07
        rect(slide, x, 1.90, 1.72, 1.10, fill=WHITE, line=accent, radius=True, width=2)
        add_text(slide, num, x + 0.12, 2.06, 0.32, 0.24, size=11, color=accent, bold=True, font=FONT_EN)
        add_text(slide, label, x + 0.42, 2.06, 1.14, 0.62, size=13, color=INK, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN, line_spacing=0.95)
        if i < len(pipeline) - 1:
            add_chevron(slide, x + 1.78, 2.26, 0.22, 0.32, accent)
    add_text(slide, "同一 base task 的六种受控变体", 0.68, 3.50, 3.50, 0.32, size=14, color=INK, bold=True)
    variants = [
        ("COMPLETE", "证据闭包完整", GREEN, GREEN_LIGHT),
        ("HIDDEN-IN-STORE", "关键证据可检索", BLUE, BLUE_LIGHT),
        ("USER-ONLY", "只能询问用户", ORANGE, ORANGE_LIGHT),
        ("STALE", "旧值需被覆盖", PURPLE, PURPLE_LIGHT),
        ("IRREDUCIBLE", "冲突无法消解", RED, RED_LIGHT),
        ("NO-MEMORY", "加入干扰但无需记忆", TEAL, TEAL_LIGHT),
    ]
    for i, (name, body, accent, fill) in enumerate(variants):
        col = i % 3
        row = i // 3
        x = 0.68 + col * 4.02
        y = 3.94 + row * 1.08
        rect(slide, x, y, 3.70, 0.84, fill=fill, line=accent, radius=True)
        add_text(slide, name, x + 0.14, y + 0.13, 1.80, 0.25, size=10, color=accent, bold=True, font=FONT_EN)
        add_text(slide, body, x + 1.77, y + 0.13, 1.75, 0.30, size=11, color=INK, bold=True, align=PP_ALIGN.RIGHT)
    rect(slide, 0.68, 6.22, 11.74, 0.47, fill=INK, line=INK, radius=True)
    add_text(slide, "核心观察不是平均 accuracy，而是 decisive memory 插入/移除后，控制动作是否发生正确翻转。", 0.96, 6.32, 11.20, 0.23, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Pilot: 60 base tasks × 6 variants; MVP: 300–500 base tasks, 1,800–3,000 episodes.")

    # 12 Metrics
    slide = prs.slides.add_slide(blank)
    add_header(slide, "09 / Evaluation", "拒绝用一个总分掩盖系统差异：四层指标分别报告", 12)
    metric_cols = [
        (0.66, "MONITOR", "状态 Macro-F1\nSufficiency AUROC\nBrier / ECE\nMissing Requirement F1\nEvidence Closure F1", TEAL, TEAL_LIGHT),
        (3.72, "CONTROL", "Admissible Action Acc.\nPolicy Regret\nPremature Execution\nUnnecessary Search / Ask\nSearch-to-Closure", BLUE, BLUE_LIGHT),
        (6.78, "OUTCOME + COST", "Task Success\nExact Tool State\nUnsupported Success\nTokens / Calls / Retrieval\nSuccess–Cost Pareto", ORANGE, ORANGE_LIGHT),
        (9.84, "COUNTERFACTUAL", "Action Flip Consistency\nIrrelevant Invariance\nStale→Current Sensitivity\nSurface Robustness\nRisk–Coverage", RED, RED_LIGHT),
    ]
    for x, title, body, accent, fill in metric_cols:
        rect(slide, x, 1.90, 2.80, 3.98, fill=fill, line=accent, radius=True, width=2)
        add_text(slide, title, x + 0.16, 2.18, 2.48, 0.32, size=13, color=accent, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
        line(slide, x + 0.24, 2.68, 2.32, 0.02, color=accent)
        add_text(slide, body, x + 0.24, 2.96, 2.32, 2.36, size=13, color=INK, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.20, font=FONT_EN)
    rect(slide, 1.56, 6.18, 10.22, 0.52, fill=INK, line=INK, radius=True)
    add_text(slide, "同样的最终成功率，可能对应完全不同的记忆控制质量与用户成本。", 1.84, 6.30, 9.66, 0.26, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Primary reporting groups: Monitor / Control / Outcome / Cost. Composite utility is secondary.")

    # 13 Baselines
    slide = prs.slides.add_slide(blank)
    add_header(slide, "10 / Baselines", "Baseline 不是越多越好，而是要能定位瓶颈", 13)
    ladder = [
        (0.72, 5.76, 11.86, "ORACLE DECOMPOSITION", "OraclePacket / OracleMonitor / OracleRetriever / OracleController / OracleAll", RED, RED_LIGHT),
        (1.25, 4.84, 10.80, "AGENTIC CONTROL", "Self-RAG-style router · MemSearcher · Memory-R1 · AgeMem · MemCon", ORANGE, ORANGE_LIGHT),
        (1.78, 3.92, 9.74, "MEMORY SYSTEM", "Mem0 / LangMem · A-MEM · StratMem-style reader · StateMem wrapper", PURPLE, PURPLE_LIGHT),
        (2.31, 3.00, 8.68, "RETRIEVAL", "BM25 · dense · hybrid · Retrieve-Once@k · FullHistory", BLUE, BLUE_LIGHT),
        (2.84, 2.08, 7.62, "SANITY", "NoMemory · AlwaysSearch · AlwaysAsk · AlwaysExecute · confidence threshold", TEAL, TEAL_LIGHT),
    ]
    for x, y, w, title, body, accent, fill in ladder:
        rect(slide, x, y, w, 0.70, fill=fill, line=accent, radius=True)
        add_text(slide, title, x + 0.20, y + 0.12, 2.05, 0.25, size=10, color=accent, bold=True, font=FONT_EN)
        add_text(slide, body, x + 2.28, y + 0.10, w - 2.48, 0.30, size=11, color=INK, bold=True, align=PP_ALIGN.RIGHT, font=FONT_EN)
    add_text(slide, "公平比较约束", 0.72, 6.74, 1.44, 0.26, size=11, color=MUTED, bold=True)
    add_text(slide, "固定 backbone、tool schema、token、搜索轮数与可见 memory source。", 2.18, 6.72, 9.80, 0.28, size=12, color=INK, bold=True)
    add_footer(slide, "MemCon is the central recent control baseline; oracles isolate retrieval, monitoring, control and execution.")

    # 14 Experiments
    slide = prs.slides.add_slide(blank)
    add_header(slide, "11 / Experiments", "实验设计：先证明 benchmark 有效，再比较方法", 14)
    experiments = [
        ("E0", "Validity", "能力梯度 / 人工一致率", TEAL, TEAL_LIGHT),
        ("E1", "Main", "模型与系统主表", BLUE, BLUE_LIGHT),
        ("E2", "Decomposition", "四种 oracle 干预", PURPLE, PURPLE_LIGHT),
        ("E3", "Counterfactual", "六变体动作翻转", ORANGE, ORANGE_LIGHT),
        ("E4", "Calibration", "risk–coverage / ECE", RED, RED_LIGHT),
        ("E5", "Budget", "memory / top-k / rounds / size", GREEN, GREEN_LIGHT),
        ("E6", "Generalization", "跨任务域迁移", TEAL, TEAL_LIGHT),
    ]
    for i, (eid, title, body, accent, fill) in enumerate(experiments):
        col = i % 4
        row = i // 4
        x = 0.68 + col * 3.04
        y = 1.90 + row * 1.24
        rect(slide, x, y, 2.76, 0.98, fill=fill, line=accent, radius=True)
        add_text(slide, eid, x + 0.16, y + 0.15, 0.46, 0.25, size=12, color=accent, bold=True, font=FONT_EN)
        add_text(slide, title, x + 0.72, y + 0.13, 1.80, 0.25, size=13, color=accent, bold=True, font=FONT_EN)
        add_text(slide, body, x + 0.16, y + 0.53, 2.42, 0.25, size=10, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "可证伪假设", 0.70, 4.62, 1.54, 0.30, size=13, color=INK, bold=True)
    hyps = [
        "H1  最终正确率会高估 memory control，因为存在猜测与不合法证据。",
        "H2  强模型在 USER_ONLY / IRREDUCIBLE 上仍会过早承诺。",
        "H3  固定 retrieve-once 帮助 hidden-in-store，却伤害 no-memory 与干扰场景。",
        "H4  MemCon 可能改善 success–cost，但不一定校准 sufficiency。",
    ]
    add_list(slide, hyps, 0.70, 5.00, 11.78, 1.55, size=12, bullet_color=RED, gap=3)
    add_footer(slide, "If simple confidence threshold approaches oracle and H1–H4 fail, stop or move to harder tasks.")

    # 15 Stage-2 method
    slide = prs.slides.add_slide(blank)
    add_header(slide, "12 / Stage 2", "Benchmark 暴露失败后，再做 MetaMem Controller", 15)
    modules = [
        (0.68, "Persistent\nStore", MUTED, GRAY_LIGHT),
        (2.55, "Evidence\nRanker", BLUE, BLUE_LIGHT),
        (4.42, "Setwise\nVerifier", PURPLE, PURPLE_LIGHT),
        (6.29, "State\nMonitor", TEAL, TEAL_LIGHT),
        (8.16, "Control\nPolicy", ORANGE, ORANGE_LIGHT),
        (10.03, "Execution\nGate", RED, RED_LIGHT),
    ]
    for i, (x, label, accent, fill) in enumerate(modules):
        rect(slide, x, 2.16, 1.55, 1.14, fill=fill, line=accent, radius=True, width=2)
        add_text(slide, label, x + 0.10, 2.42, 1.35, 0.54, size=14, color=accent, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN, line_spacing=0.92)
        if i < len(modules) - 1:
            add_chevron(slide, x + 1.61, 2.54, 0.21, 0.34, accent)
    actions = [
        ("SEARCH", BLUE, BLUE_LIGHT),
        ("ASK", ORANGE, ORANGE_LIGHT),
        ("EXECUTE", GREEN, GREEN_LIGHT),
        ("ABSTAIN", RED, RED_LIGHT),
    ]
    for i, (name, accent, fill) in enumerate(actions):
        x = 3.10 + i * 1.82
        rect(slide, x, 3.92, 1.52, 0.58, fill=fill, line=accent, radius=True)
        add_text(slide, name, x + 0.08, 4.08, 1.36, 0.24, size=11, color=accent, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(slide, "控制动作", 1.40, 4.08, 1.38, 0.25, size=12, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    add_chevron(slide, 2.82, 4.03, 0.22, 0.32, ORANGE)
    rect(slide, 0.68, 5.10, 11.76, 1.24, fill=INK, line=INK, radius=True)
    add_text(slide, "设计原则", 0.98, 5.36, 1.15, 0.28, size=13, color=TEAL_LIGHT, bold=True)
    add_text(slide, "只有 verifier 通过，或 policy 明确接受风险时，Execution Gate 才允许最终行动。", 2.24, 5.30, 9.80, 0.35, size=17, color=WHITE, bold=True)
    add_text(slide, "目标：把“看见相关记忆”升级为“形成最小充分证据闭包后再行动”。", 2.24, 5.79, 9.80, 0.28, size=13, color=WHITE)
    add_footer(slide, "Method is intentionally second-stage: benchmark first, intervention second.")

    # 16 Learning and reward
    slide = prs.slides.add_slide(blank)
    add_header(slide, "12 / Stage 2", "排序学习放在哪里：逐条判断、相对选择、集合闭包", 16)
    add_card(slide, 0.68, 1.90, 3.68, 1.72, "Pointwise", "预测单条 memory 的相关性、时效性、authority 与 slot coverage。", accent=TEAL, fill=TEAL_LIGHT, title_size=18, body_size=13)
    add_card(slide, 4.82, 1.90, 3.68, 1.72, "Pairwise", "在 decisive evidence 与 hard negative 之间学习相对效用，适合下一条证据 / 动作选择。", accent=BLUE, fill=BLUE_LIGHT, title_size=18, body_size=13)
    add_card(slide, 8.96, 1.90, 3.68, 1.72, "Listwise / Setwise", "选择能闭合全部需求的最小证据集合，处理“每条都相关但组合仍不充分”。", accent=PURPLE, fill=PURPLE_LIGHT, title_size=18, body_size=13)
    add_text(slide, "训练路线", 0.70, 4.04, 1.20, 0.30, size=13, color=INK, bold=True)
    steps = [
        ("1", "Rule / prompt", "先验证 benchmark", TEAL),
        ("2", "SFT", "gold state / action / evidence", BLUE),
        ("3", "Ranking", "pairwise + setwise closure", PURPLE),
        ("4", "Scale", "0.6B smoke → 7B / 8B", ORANGE),
        ("5", "Optional RL", "仅处理停止与预算", RED),
    ]
    for i, (num, title, body, accent) in enumerate(steps):
        x = 0.70 + i * 2.42
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(4.50), Inches(0.46), Inches(0.46))
        circ.fill.solid(); circ.fill.fore_color.rgb = rgb(accent); circ.line.color.rgb = rgb(accent)
        add_text(slide, num, x, 4.60, 0.46, 0.20, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(slide, title, x + 0.56, 4.49, 1.50, 0.26, size=11, color=accent, bold=True, font=FONT_EN)
        add_text(slide, body, x + 0.56, 4.80, 1.62, 0.42, size=9, color=INK, bold=True)
    rect(slide, 0.70, 5.70, 11.94, 0.84, fill=WHITE, line=RED, radius=True, width=2)
    add_text(slide, "最终 reward", 0.98, 5.98, 1.22, 0.25, size=12, color=RED, bold=True)
    add_text(slide, "Exact task success + admissible action + evidence closure − search/clarify cost − premature/unsupported action", 2.18, 5.93, 9.98, 0.34, size=13, color=INK, bold=True, font=FONT_EN)
    add_text(slide, "第一版不训练独立 GRM：结构化 gold 与环境执行已提供 RLVR 风格信号。", 1.02, 6.69, 11.30, 0.26, size=12, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Ranking score shapes evidence selection; it does not replace verifiable outcome reward.")

    # 17 MVP and gates
    slide = prs.slides.add_slide(blank)
    add_header(slide, "13 / Feasibility", "最小闭环：1–2 周判断这个问题是否真的成立", 17)
    add_number_card(slide, 0.70, 1.90, 2.22, 1.42, "20 × 6", "BASE TASKS × VARIANTS", TEAL, "120 episodes")
    add_number_card(slide, 3.18, 1.90, 2.22, 1.42, "1", "DETERMINISTIC TOOL", BLUE, "calendar / travel")
    add_number_card(slide, 5.66, 1.90, 2.22, 1.42, "4", "SANITY BASELINES", ORANGE, "Execute / Retrieve / Full / Oracle")
    add_number_card(slide, 8.14, 1.90, 2.22, 1.42, "1–2", "WEEKS", PURPLE, "without training")
    add_number_card(slide, 10.62, 1.90, 2.02, 1.42, "4×4090", "AVAILABLE", RED, "7B/8B LoRA-ready")
    add_text(slide, "Go / No-Go Gate", 0.72, 3.80, 2.20, 0.32, size=14, color=INK, bold=True, font=FONT_EN)
    gates = [
        "OracleAll 接近可解上界，NoMemory 与 OraclePacket 有显著差距。",
        "至少两个 baseline 在 premature execution / unnecessary search 上呈现不同 trade-off。",
        "强模型的 counterfactual action flip 不能接近满分；FullHistory 不能轻易饱和。",
        "人工 state / admissible action 一致率足够高，且 ≥80% 主指标可规则化计算。",
    ]
    add_list(slide, gates, 0.72, 4.24, 11.74, 1.66, size=13, bullet_color=TEAL, gap=5)
    rect(slide, 0.72, 6.22, 11.72, 0.50, fill=INK, line=INK, radius=True)
    add_text(slide, "GPU 不是核心壁垒：主要资产是 generator、simulator、evaluator 与 baseline harness。", 0.98, 6.33, 11.20, 0.25, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Compute plan: CPU deterministic environments; 0.6B smoke; 7B/8B inference and light LoRA/SFT.")

    # 18 Roadmap
    slide = prs.slides.add_slide(blank)
    add_header(slide, "14 / Roadmap", "10 周执行路线：每一阶段都有可停止的检查点", 18)
    roadmap = [
        ("W1", "Spec", "taxonomy + 20 samples", TEAL),
        ("W2", "Generator", "120-episode pilot", BLUE),
        ("W3", "Track A", "sanity baselines + gate", PURPLE),
        ("W4–5", "Tracks B/C", "API + 2 tool domains", ORANGE),
        ("W6", "Baselines", "retrieval + memory systems", GREEN),
        ("W7", "Controllers", "MemCon + oracle diagnosis", RED),
        ("W8", "Scale", "1,800+ episodes", TEAL),
        ("W9", "Analysis", "calibration / CF / budget", BLUE),
        ("W10", "Release", "paper + code + data card", INK),
    ]
    for i, (week, title, body, accent) in enumerate(roadmap):
        row = i % 5
        col = i // 5
        x = 0.72 + col * 6.08
        y = 1.88 + row * 0.95
        rect(slide, x, y, 5.72, 0.74, fill=WHITE, line=accent, radius=True)
        rect(slide, x, y, 0.86, 0.74, fill=accent, line=accent, radius=True)
        add_text(slide, week, x + 0.05, y + 0.23, 0.76, 0.24, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(slide, title, x + 1.03, y + 0.12, 1.24, 0.25, size=12, color=accent, bold=True, font=FONT_EN)
        add_text(slide, body, x + 2.22, y + 0.12, 3.24, 0.30, size=10, color=INK, bold=True, align=PP_ALIGN.RIGHT, font=FONT_EN)
    rect(slide, 0.72, 6.74, 12.00, 0.25, fill=TEAL_LIGHT, line=TEAL_LIGHT)
    add_text(slide, "关键决策点：W3 pilot gate → W7 method need → W9 是否支撑“通用 agentic memory control”论点。", 0.92, 6.75, 11.58, 0.20, size=10, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Roadmap prioritizes benchmark validity before dataset scale and method complexity.")

    # 19 Contributions and boundaries
    slide = prs.slides.add_slide(blank)
    add_header(slide, "15 / Contribution", "预期贡献与论文主线", 19)
    contribs = [
        ("01", "Problem", "定义 memory monitoring and control before commitment。", TEAL, TEAL_LIGHT),
        ("02", "Benchmark", "gold state、evidence closure、admissible action 与 executable outcome。", BLUE, BLUE_LIGHT),
        ("03", "Protocol", "Monitor / Control / Outcome / Cost + oracle decomposition。", ORANGE, ORANGE_LIGHT),
        ("04", "Evidence", "系统比较 long-context、RAG、memory backend 与 agentic controllers。", RED, RED_LIGHT),
    ]
    for i, (num, label, body, accent, fill) in enumerate(contribs):
        y = 1.86 + i * 1.05
        rect(slide, 0.70, y, 7.58, 0.80, fill=fill, line=accent, radius=True)
        add_text(slide, num, 0.92, y + 0.23, 0.46, 0.24, size=11, color=accent, bold=True, font=FONT_EN)
        add_text(slide, label, 1.50, y + 0.19, 1.16, 0.28, size=13, color=accent, bold=True, font=FONT_EN)
        add_text(slide, body, 2.70, y + 0.17, 5.28, 0.32, size=12, color=INK, bold=True)
    rect(slide, 8.62, 1.86, 4.02, 4.87, fill=INK, line=INK, radius=True)
    add_text(slide, "Take-home", 8.98, 2.24, 3.30, 0.34, size=14, color=TEAL_LIGHT, bold=True, font=FONT_EN)
    add_text(slide, "Agentic memory 的关键不只是“能存、能取”，而是：", 8.98, 2.88, 3.16, 0.78, size=18, color=WHITE, bold=True)
    add_text(slide, "知道何时已经足够，\n何时还应继续找，\n何时必须向人询问，\n何时应该停止承诺。", 8.98, 3.84, 3.16, 1.82, size=20, color=WHITE, bold=True, line_spacing=1.12)
    add_badge(slide, "BENCHMARK FIRST", 9.18, 6.00, 2.90, fill=TEAL_LIGHT, color=TEAL, size=11)
    add_footer(slide, "Recommended title: MetaMemBench: Evaluating Monitoring and Control of Agent Memory under Uncertainty.")

    # 20 References
    slide = prs.slides.add_slide(blank)
    add_header(slide, "References", "核心参考文献", 20, "完整调研与问题推导见配套 proposal 文档")
    refs_left = [
        (1, "Memory in the Age of AI Agents", "Survey · arXiv 2512.13564", "https://arxiv.org/abs/2512.13564"),
        (2, "StratMem-Bench", "ACL 2026 Main · Long Paper", "https://aclanthology.org/2026.acl-long.1491/"),
        (3, "Mem2ActBench", "ACL 2026 Main · Long Paper", "https://aclanthology.org/2026.acl-long.370/"),
        (4, "MemoryArena", "arXiv 2602.16313", "https://arxiv.org/abs/2602.16313"),
        (5, "AMemGym", "ICLR 2026", "https://arxiv.org/abs/2603.01966"),
        (6, "MemOps", "arXiv 2607.12893", "https://arxiv.org/abs/2607.12893"),
        (7, "MemTrace", "arXiv 2605.28732", "https://arxiv.org/abs/2605.28732"),
    ]
    refs_right = [
        (8, "TANGLE", "arXiv 2608.13921", "https://arxiv.org/abs/2608.13921"),
        (9, "StateMemBench", "arXiv 2608.19652", "https://arxiv.org/abs/2608.19652"),
        (10, "AgeMem", "ACL 2026 Main · Long Paper", "https://aclanthology.org/2026.acl-long.981/"),
        (11, "Memory as a Controlled Process (MemCon)", "arXiv 2607.13591", "https://arxiv.org/abs/2607.13591"),
        (12, "StreamMemBench", "arXiv 2606.14571", "https://arxiv.org/abs/2606.14571"),
        (13, "CIMemories", "ICLR 2026", "https://arxiv.org/abs/2511.14937"),
        (14, "Decision-Aware Memory Cards / CICL", "arXiv 2606.08151", "https://arxiv.org/abs/2606.08151"),
    ]
    for i, title, venue, url in refs_left:
        add_reference(slide, i, title, venue, url, 0.72, 1.86 + (i - 1) * 0.69, 5.86)
    for row, (i, title, venue, url) in enumerate(refs_right):
        add_reference(slide, i, title, venue, url, 6.78, 1.86 + row * 0.69, 5.86)
    rect(slide, 0.72, 6.82, 11.92, 0.20, fill=TEAL_LIGHT, line=TEAL_LIGHT)
    add_text(slide, "Proposal source: docs/ideas/2026-09-03-metamembench-agentic-memory-control-proposal.md", 0.86, 6.82, 11.56, 0.20, size=8, color=TEAL, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"saved: {OUT}")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_deck()
